from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from pymongo import MongoClient
from flask_restful import Api
from flask_swagger import swagger
import bcrypt
import json
import os
from dotenv import load_dotenv
import requests
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
import logging
import sys
from datetime import datetime, timedelta
from io import BytesIO
import google.generativeai as genai
from PIL import Image
from pptx import Presentation
import re
from functools import wraps
from collections import defaultdict

# Configure logging
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s [%(levelname)s] %(message)s',
    handlers=[
        logging.FileHandler('app.log'),
        logging.StreamHandler(sys.stdout)
    ]
)

# Create logger instance
logger = logging.getLogger(__name__)

# Load environment variables from .env file
load_dotenv()

# Configure upload folder from environment variables
UPLOAD_FOLDER = os.path.join(os.getcwd(), os.getenv('UPLOAD_FOLDER', 'uploaded_slides'))
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
ALLOWED_EXTENSIONS = {"pdf", "txt", "docx", "pptx", "png", "jpg", "jpeg"}

app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

# Configure CORS
CORS(app, resources={
    r"/*": {
        "origins": os.getenv('CORS_ORIGINS', 'http://localhost:3000').split(',')
    }
})

# Configuration from environment variables
app.config['DEBUG'] = os.getenv('FLASK_DEBUG', 'True').lower() == 'true'
app.config['ENV'] = os.getenv('FLASK_ENV', 'development')
app.secret_key = os.getenv('SECRET_KEY', 'fallback-secret-key-change-in-production')

# MongoDB configuration from environment variables
mongodb_uri = os.getenv('MONGODB_URI', 'mongodb://localhost:27017/')
database_name = os.getenv('DATABASE_NAME', 'eduplanner')

# Connect to MongoDB
client = MongoClient(mongodb_uri)
db = client[database_name]  # Database Name
users_collection = db.users  # Collection Name

# Rate limiting storage
login_attempts = defaultdict(list)
RATE_LIMIT_WINDOW = 300  # 5 minutes in seconds
MAX_ATTEMPTS_PER_WINDOW = 5

# Add CORS headers to all responses
@app.after_request
def after_request(response):
    response.headers['Access-Control-Allow-Origin'] = 'http://localhost:3000'
    response.headers['Access-Control-Allow-Methods'] = 'GET, POST, PUT, DELETE, OPTIONS'
    response.headers['Access-Control-Allow-Headers'] = 'Content-Type'
    response.headers['Access-Control-Max-Age'] = '3600'
    return response

# Handle OPTIONS preflight requests
@app.route('/slides/<filename>', methods=['OPTIONS'])
def handle_preflight(filename):
    response = jsonify({'message': 'OK'})
    return response, 200

# Password validation function
def validate_password(password):
    """Validate password strength"""
    if len(password) < 8:
        return False, "Password must be at least 8 characters long"
    
    if not re.search(r'[A-Z]', password):
        return False, "Password must contain at least one uppercase letter"
    
    if not re.search(r'[a-z]', password):
        return False, "Password must contain at least one lowercase letter"
    
    if not re.search(r'\d', password):
        return False, "Password must contain at least one number"
    
    if not re.search(r'[!@#$%^&*()_+\-=\[\]{};:"\\|,.<>\/?]', password):
        return False, "Password must contain at least one special character"
    
    return True, "Password is valid"

# Rate limiting decorator
def rate_limit(max_attempts=MAX_ATTEMPTS_PER_WINDOW, window=RATE_LIMIT_WINDOW):
    def decorator(f):
        @wraps(f)
        def decorated_function(*args, **kwargs):
            client_ip = request.environ.get('HTTP_X_FORWARDED_FOR', request.environ.get('REMOTE_ADDR', 'unknown'))
            current_time = datetime.now()
            
            # Clean old attempts
            login_attempts[client_ip] = [
                attempt_time for attempt_time in login_attempts[client_ip]
                if (current_time - attempt_time).total_seconds() < window
            ]
            
            # Check if rate limit exceeded
            if len(login_attempts[client_ip]) >= max_attempts:
                logger.warning(f"Rate limit exceeded for IP: {client_ip}")
                return jsonify({
                    "status": "error",
                    "message": "Too many login attempts. Please try again later."
                }), 429
            
            # Record this attempt
            login_attempts[client_ip].append(current_time)
            
            return f(*args, **kwargs)
        return decorated_function
    return decorator

# Input validation function
def validate_input(data, required_fields):
    """Validate input data"""
    errors = []
    
    for field in required_fields:
        if field not in data or not data[field] or not data[field].strip():
            errors.append(f"{field} is required")
        elif len(data[field].strip()) > 255:  # Prevent extremely long inputs
            errors.append(f"{field} is too long (max 255 characters)")
    
    return errors

# Signup Route
@app.route("/signup", methods=["POST"])
def signup():
    try:
        data = request.json
        if not data:
            return jsonify({"error": "No data provided"}), 400
        
        # Validate required fields
        required_fields = ["firstname", "lastname", "studentId", "email", "password"]
        validation_errors = validate_input(data, required_fields)
        
        # Handle both 'name' field (from frontend) and separate firstname/lastname
        if 'name' in data and not data.get('firstname'):
            # Split full name into first and last
            name_parts = data['name'].strip().split()
            data['firstname'] = name_parts[0] if name_parts else ''
            data['lastname'] = ' '.join(name_parts[1:]) if len(name_parts) > 1 else ''
        
        if validation_errors:
            return jsonify({"error": "; ".join(validation_errors)}), 400
        
        firstname = data.get("firstname").strip()
        lastname = data.get("lastname").strip()
        studentId = data.get("studentId").strip()
        email = data.get("email").strip().lower()
        password = data.get("password")
        
        # Validate email format
        email_pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
        if not re.match(email_pattern, email):
            return jsonify({"error": "Invalid email format"}), 400
        
        # Validate password strength
        is_valid, password_message = validate_password(password)
        if not is_valid:
            return jsonify({"error": password_message}), 400
        
        # Check if email already exists
        if users_collection.find_one({"email": email}):
            return jsonify({"error": "Email already exists"}), 400
        
        # Check if student ID already exists
        if users_collection.find_one({"studentId": studentId}):
            return jsonify({"error": "Student ID already exists"}), 400
        
        # Hash password
        hashed_pw = bcrypt.hashpw(password.encode("utf-8"), bcrypt.gensalt())
        
        # Create user
        user_data = {
            "firstname": firstname,
            "lastname": lastname,
            "studentId": studentId,
            "email": email,
            "password": hashed_pw,
            "created_at": datetime.now(),
            "is_active": True
        }
        
        result = users_collection.insert_one(user_data)
        
        logger.info(f"User created successfully: {email}")
        
        return jsonify({
            "status": "success",
            "message": "User created successfully",
            "user_id": str(result.inserted_id)
        }), 201
        
    except Exception as e:
        logger.error(f"Signup error: {str(e)}")
        return jsonify({"error": "Internal server error"}), 500

# Login Route with Rate Limiting
@app.route("/login", methods=["POST"])
@rate_limit()
def login():
    try:
        data = request.get_json()
        if not data:
            return jsonify({"status": "error", "message": "No data provided"}), 400
        
        # Validate required fields
        required_fields = ["email", "password"]
        validation_errors = validate_input(data, required_fields)
        
        if validation_errors:
            return jsonify({"status": "error", "message": "; ".join(validation_errors)}), 400
        
        email = data.get("email").strip().lower()
        password = data.get("password")
        
        # Validate email format
        email_pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
        if not re.match(email_pattern, email):
            return jsonify({"status": "error", "message": "Invalid email format"}), 400
        
        # Find user
        user = users_collection.find_one({"email": email})
        
        if user and bcrypt.checkpw(password.encode("utf-8"), user["password"]):
            # Check if user account is active
            if not user.get("is_active", True):
                return jsonify({
                    "status": "error",
                    "message": "Account is deactivated. Please contact support."
                }), 403
            
            # Update last login
            users_collection.update_one(
                {"_id": user["_id"]},
                {"$set": {"last_login": datetime.now()}}
            )
            
            logger.info(f"Successful login for user: {email}")
            
            return jsonify({
                "status": "success",
                "message": "Logged in successfully",
                "user": {
                    "id": str(user["_id"]),
                    "firstname": user.get("firstname"),
                    "lastname": user.get("lastname"),
                    "email": user.get("email"),
                    "studentId": user.get("studentId")
                }
            }), 200
        else:
            logger.warning(f"Failed login attempt for email: {email}")
            return jsonify({
                "status": "fail",
                "message": "Invalid credentials"
            }), 401
            
    except Exception as e:
        logger.error(f"Login error: {str(e)}")
        return jsonify({"status": "error", "message": "Internal server error"}), 500

# Configure Gemini API from environment variables
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
if not GEMINI_API_KEY:
    logger.error("GEMINI_API_KEY not found in environment variables!")
    raise ValueError("GEMINI_API_KEY environment variable is required")

genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel('gemini-1.5-flash')

# Upload configuration
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pdf', 'png', 'jpg', 'jpeg', 'ppt', 'pptx'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Create upload directory if it doesn't exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
    return text

def extract_text_from_image(file_path):
    """Extract text from image using OCR"""
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        print(f"Error extracting text from image: {e}")
        return ""

def extract_text_from_ppt(file_path):
    """Extract text from PowerPoint file"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from PowerPoint: {e}")
        return ""

def generate_summary(text):
    """Generate summary using Gemini API"""
    try:
        prompt = f"""
        Please provide a comprehensive summary of the following content from slides/presentation:
        
        {text}
        
        Please structure the summary with:
        1. Main topic/title
        2. Key points (bullet points)
        3. Important details
        4. Conclusion
        
        Make it clear and concise while covering all important information.
        """
        
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating summary: {e}")
        return "Error generating summary. Please try again."

def generate_quiz(text):
    """Generate quiz questions using Gemini API"""
    try:
        prompt = f"""
        Based on the following content from slides/presentation, create a quiz with 5 multiple-choice questions:
        
        {text}
        
        Please format the response as JSON with the following structure:
        {{
            "questions": [
                {{
                    "question": "Question text here?",
                    "options": ["A) Option 1", "B) Option 2", "C) Option 3", "D) Option 4"],
                    "correct_answer": "A"
                }}
            ]
        }}
        
        Make sure the questions test understanding of key concepts from the content.
        """
        
        response = model.generate_content(prompt)
        # Try to parse JSON from response
        response_text = response.text.strip()
        # Remove markdown code blocks if present
        if response_text.startswith('```json'):
            response_text = response_text[7:-3]
        elif response_text.startswith('```'):
            response_text = response_text[3:-3]
        
        quiz_data = json.loads(response_text)
        return quiz_data
    except Exception as e:
        print(f"Error generating quiz: {e}")
        # Return a fallback quiz structure
        return {
            "questions": [
                {
                    "question": "Error generating quiz questions. Please try uploading the file again.",
                    "options": ["A) Try again", "B) Try again", "C) Try again", "D) Try again"],
                    "correct_answer": "A"
                }
            ]
        }

@app.route('/')
def index():
    return jsonify({"message": "Slide Analyzer API is running!"})

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        # Extract text based on file type
        file_extension = filename.rsplit('.', 1)[1].lower()
        
        if file_extension == 'pdf':
            extracted_text = extract_text_from_pdf(file_path)
        elif file_extension in ['png', 'jpg', 'jpeg']:
            extracted_text = extract_text_from_image(file_path)
        elif file_extension in ['ppt', 'pptx']:
            extracted_text = extract_text_from_ppt(file_path)
        else:
            return jsonify({'error': 'Unsupported file type'}), 400
        
        if not extracted_text.strip():
            return jsonify({'error': 'No text could be extracted from the file'}), 400
        
        return jsonify({
            'message': 'File uploaded successfully',
            'filename': filename,
            'extracted_text': extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text
        })
    
    return jsonify({'error': 'Invalid file type'}), 400

@app.route('/generate-summary', methods=['POST'])
def generate_summary_endpoint():
    data = request.get_json()
    
    if not data or 'filename' not in data:
        return jsonify({'error': 'Filename required'}), 400
    
    filename = data['filename']
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    
    if not os.path.exists(file_path):
        return jsonify({'error': 'File not found'}), 404
    
    # Extract text based on file type
    file_extension = filename.rsplit('.', 1)[1].lower()
    
    if file_extension == 'pdf':
        extracted_text = extract_text_from_pdf(file_path)
    elif file_extension in ['png', 'jpg', 'jpeg']:
        extracted_text = extract_text_from_image(file_path)
    elif file_extension in ['ppt', 'pptx']:
        extracted_text = extract_text_from_ppt(file_path)
    else:
        return jsonify({'error': 'Unsupported file type'}), 400
    
    if not extracted_text.strip():
        return jsonify({'error': 'No text could be extracted from the file'}), 400
    
    summary = generate_summary(extracted_text)
    
    return jsonify({
        'summary': summary,
        'filename': filename
    })

@app.route('/generate-quiz', methods=['POST'])
def generate_quiz_endpoint():
    data = request.get_json()
    
    if not data or 'filename' not in data:
        return jsonify({'error': 'Filename required'}), 400
    
    filename = data['filename']
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    
    if not os.path.exists(file_path):
        return jsonify({'error': 'File not found'}), 404
    
    # Extract text based on file type
    file_extension = filename.rsplit('.', 1)[1].lower()
    
    if file_extension == 'pdf':
        extracted_text = extract_text_from_pdf(file_path)
    elif file_extension in ['png', 'jpg', 'jpeg']:
        extracted_text = extract_text_from_image(file_path)
    elif file_extension in ['ppt', 'pptx']:
        extracted_text = extract_text_from_ppt(file_path)
    else:
        return jsonify({'error': 'Unsupported file type'}), 400
    
    if not extracted_text.strip():
        return jsonify({'error': 'No text could be extracted from the file'}), 400
    
    quiz_data = generate_quiz(extracted_text)
    
    return jsonify({
        'quiz': quiz_data,
        'filename': filename
    })

if __name__ == '__main__':
    app.run(debug=True, port=5000)
