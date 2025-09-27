"""
File management module for EduMaster application.
This module handles file uploads, processing, and text extraction from various file types.
"""

import os
import logging
from datetime import datetime
from flask import request, jsonify, send_from_directory
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
from PIL import Image
from pptx import Presentation

from database import files_collection
from utils import allowed_file, validate_input, format_file_size, log_user_action

logger = logging.getLogger(__name__)

# File configuration
ALLOWED_EXTENSIONS = {'pdf', 'png', 'jpg', 'jpeg', 'ppt', 'pptx', 'txt', 'docx'}
MAX_FILE_SIZE = 16 * 1024 * 1024  # 16MB

class FileProcessor:
    """Handles file processing and text extraction"""
    
    @staticmethod
    def extract_text_from_pdf(file_path):
        """Extract text from PDF file"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
        return text

    @staticmethod
    def extract_text_from_image(file_path):
        """Extract text from image using OCR"""
        try:
            # Note: This requires pytesseract to be installed and configured
            import pytesseract
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text
        except ImportError:
            logger.error("pytesseract not available for OCR")
            return "OCR not available - pytesseract not installed"
        except Exception as e:
            logger.error(f"Error extracting text from image: {e}")
            return ""

    @staticmethod
    def extract_text_from_ppt(file_path):
        """Extract text from PowerPoint file"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PowerPoint: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_txt(file_path):
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encoding if UTF-8 fails
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                logger.error(f"Error reading text file with latin-1 encoding: {e}")
                return ""
        except Exception as e:
            logger.error(f"Error extracting text from text file: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_docx(file_path):
        """Extract text from Word document"""
        try:
            # Try importing docx
            try:
                from docx import Document
                doc = Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text
            except ImportError:
                logger.error("python-docx not available for DOCX text extraction")
                return "DOCX text extraction not available - python-docx not installed"
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {e}")
            return ""

    @classmethod
    def extract_text_by_extension(cls, file_path, file_extension):
        """Extract text based on file extension"""
        if file_extension == 'pdf':
            return cls.extract_text_from_pdf(file_path)
        elif file_extension in ['png', 'jpg', 'jpeg']:
            return cls.extract_text_from_image(file_path)
        elif file_extension in ['ppt', 'pptx']:
            return cls.extract_text_from_ppt(file_path)
        elif file_extension == 'txt':
            return cls.extract_text_from_txt(file_path)
        elif file_extension == 'docx':
            return cls.extract_text_from_docx(file_path)
        else:
            return ""

def setup_upload_folder():
    """Create upload folder if it doesn't exist"""
    upload_folder = os.path.join(os.getcwd(), os.getenv('UPLOAD_FOLDER', 'uploads'))
    if not os.path.exists(upload_folder):
        os.makedirs(upload_folder)
    return upload_folder

def register_file_routes(app):
    """Register file management routes with the Flask app"""
    
    # Setup upload folder and configure app
    upload_folder = setup_upload_folder()
    app.config['UPLOAD_FOLDER'] = upload_folder
    app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE
    
    @app.route('/upload', methods=['POST'])
    def upload_file():
        """Upload and process file endpoint"""
        try:
            if 'file' not in request.files:
                return jsonify({'error': 'No file uploaded'}), 400
            
            file = request.files['file']
            if file.filename == '':
                return jsonify({'error': 'No file selected'}), 400
            
            # Get additional data from form
            user_id = request.form.get('user_id', 'anonymous')
            course_id = request.form.get('course_id', 'general')
            
            if file and allowed_file(file.filename, ALLOWED_EXTENSIONS):
                filename = secure_filename(file.filename)
                
                # Generate unique filename to prevent conflicts
                base_name, extension = os.path.splitext(filename)
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                unique_filename = f"{base_name}_{timestamp}{extension}"
                
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                file.save(file_path)
                
                # Get file info
                file_size = os.path.getsize(file_path)
                file_extension = filename.rsplit('.', 1)[1].lower()
                
                # Extract text based on file type
                extracted_text = FileProcessor.extract_text_by_extension(file_path, file_extension)
                
                if not extracted_text.strip():
                    # Don't fail completely, but warn user
                    extracted_text = "No text could be extracted from this file."
                    logger.warning(f"No text extracted from file: {filename}")
                
                # Store file metadata in database
                file_metadata = {
                    'filename': unique_filename,
                    'original_filename': file.filename,
                    'file_path': file_path,
                    'file_size': file_size,
                    'file_extension': file_extension,
                    'extracted_text': extracted_text,
                    'user_id': user_id,
                    'course_id': course_id or 'general',
                    'upload_date': datetime.now(),
                    'is_active': True
                }
                
                result = files_collection.insert_one(file_metadata)
                file_id = str(result.inserted_id)
                
                # Log successful upload
                log_user_action(user_id, "file_upload", {
                    "filename": file.filename,
                    "file_size": file_size,
                    "file_type": file_extension
                })
                
                logger.info(f"File uploaded successfully: {file.filename} by user {user_id}")
                
                return jsonify({
                    'status': 'success',
                    'message': 'File uploaded successfully',
                    'file_id': file_id,
                    'filename': unique_filename,
                    'original_filename': file.filename,
                    'file_size': format_file_size(file_size),
                    'extracted_text_preview': extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text
                }), 201
            
            return jsonify({'error': 'Invalid file type. Allowed types: PDF, PPT, PPTX, PNG, JPG, JPEG'}), 400
            
        except Exception as e:
            logger.error(f"Error uploading file: {str(e)}")
            return jsonify({'error': 'Internal server error during file upload'}), 500

    @app.route('/course-files/<course_id>', methods=['GET'])
    def get_course_files(course_id):
        """Get files for a specific course"""
        try:
            user_id = request.args.get('user_id')
            if not user_id:
                return jsonify({'error': 'User ID is required'}), 400
            
            # Build query
            query = {
                'course_id': course_id,
                'user_id': user_id,
                'is_active': True
            }
            
            files = files_collection.find(query).sort('upload_date', -1)
            
            file_list = []
            for file in files:
                file_data = {
                    'file_id': str(file['_id']),
                    'filename': file['filename'],
                    'original_filename': file['original_filename'],
                    'file_size': file['file_size'],
                    'file_size_formatted': format_file_size(file['file_size']),
                    'file_extension': file['file_extension'],
                    'upload_date': file['upload_date'].isoformat(),
                    'extracted_text_preview': file['extracted_text'][:200] + "..." if len(file['extracted_text']) > 200 else file['extracted_text']
                }
                file_list.append(file_data)
            
            return jsonify({
                'status': 'success',
                'files': file_list,
                'total_files': len(file_list)
            })
            
        except Exception as e:
            logger.error(f"Error retrieving course files: {str(e)}")
            return jsonify({'error': 'Error retrieving files'}), 500

    @app.route('/user-files/<user_id>', methods=['GET'])
    def get_user_files(user_id):
        """Get all files for a specific user"""
        try:
            # Get query parameters
            course_id = request.args.get('course_id')
            limit = int(request.args.get('limit', 50))
            
            # Build query
            query = {
                'user_id': user_id,
                'is_active': True
            }
            
            if course_id and course_id != 'all':
                query['course_id'] = course_id
            
            files = files_collection.find(query).sort('upload_date', -1).limit(limit)
            
            file_list = []
            for file in files:
                file_data = {
                    'file_id': str(file['_id']),
                    'filename': file['filename'],
                    'original_filename': file['original_filename'],
                    'file_size': file['file_size'],
                    'file_size_formatted': format_file_size(file['file_size']),
                    'file_extension': file['file_extension'],
                    'course_id': file.get('course_id', 'general'),
                    'upload_date': file['upload_date'].isoformat(),
                    'extracted_text_preview': file['extracted_text'][:200] + "..." if len(file['extracted_text']) > 200 else file['extracted_text']
                }
                file_list.append(file_data)
            
            return jsonify({
                'status': 'success',
                'files': file_list,
                'total_files': len(file_list)
            })
            
        except Exception as e:
            logger.error(f"Error retrieving user files: {str(e)}")
            return jsonify({'error': 'Error retrieving files'}), 500

    @app.route('/file-content/<file_id>', methods=['GET'])
    def get_file_content(file_id):
        """Get specific file content and metadata"""
        try:
            from bson import ObjectId
            
            try:
                file_object_id = ObjectId(file_id)
            except:
                return jsonify({'error': 'Invalid file ID format'}), 400
            
            file_data = files_collection.find_one({'_id': file_object_id})
            
            if not file_data:
                return jsonify({'error': 'File not found'}), 404
            
            return jsonify({
                'status': 'success',
                'file': {
                    'file_id': str(file_data['_id']),
                    'filename': file_data['filename'],
                    'original_filename': file_data['original_filename'],
                    'file_size': file_data['file_size'],
                    'file_size_formatted': format_file_size(file_data['file_size']),
                    'file_extension': file_data['file_extension'],
                    'course_id': file_data.get('course_id', 'general'),
                    'extracted_text': file_data['extracted_text'],
                    'upload_date': file_data['upload_date'].isoformat()
                }
            })
            
        except Exception as e:
            logger.error(f"Error retrieving file content: {str(e)}")
            return jsonify({'error': 'Error retrieving file content'}), 500

    @app.route('/delete-file/<file_id>', methods=['DELETE'])
    def delete_file(file_id):
        """Delete a file"""
        try:
            from bson import ObjectId
            
            try:
                file_object_id = ObjectId(file_id)
            except:
                return jsonify({'error': 'Invalid file ID format'}), 400
            
            # Get file data first
            file_data = files_collection.find_one({'_id': file_object_id})
            if not file_data:
                return jsonify({'error': 'File not found'}), 404
            
            user_id = request.args.get('user_id')
            if not user_id or user_id != file_data.get('user_id'):
                return jsonify({'error': 'Unauthorized to delete this file'}), 403
            
            # Mark file as inactive (soft delete)
            result = files_collection.update_one(
                {'_id': file_object_id},
                {'$set': {'is_active': False, 'deleted_at': datetime.now()}}
            )
            
            if result.modified_count > 0:
                # Optionally, delete physical file
                try:
                    if os.path.exists(file_data['file_path']):
                        os.remove(file_data['file_path'])
                except Exception as e:
                    logger.warning(f"Could not delete physical file: {e}")
                
                log_user_action(user_id, "file_delete", {
                    "filename": file_data['original_filename'],
                    "file_id": file_id
                })
                
                return jsonify({
                    'status': 'success',
                    'message': 'File deleted successfully'
                })
            else:
                return jsonify({'error': 'Failed to delete file'}), 500
                
        except Exception as e:
            logger.error(f"Error deleting file: {str(e)}")
            return jsonify({'error': 'Error deleting file'}), 500

    @app.route('/serve-file/<filename>', methods=['GET'])
    def serve_file(filename):
        """Serve uploaded files directly for viewing"""
        try:
            # Security check: only serve files from upload directory
            upload_folder = app.config.get('UPLOAD_FOLDER', 'uploads')
            
            # Ensure we have an absolute path
            if not os.path.isabs(upload_folder):
                upload_folder = os.path.abspath(upload_folder)
            
            file_path = os.path.join(upload_folder, filename)
            
            # Security check: prevent directory traversal
            # Ensure the file path is within the upload directory
            try:
                # Resolve any relative path components
                file_path = os.path.abspath(file_path)
                upload_folder = os.path.abspath(upload_folder)
                
                # Check if the file path starts with the upload folder path
                if not file_path.startswith(upload_folder):
                    logger.warning(f"Security violation: Attempted to access file outside upload directory: {filename}")
                    return jsonify({'error': 'File not found'}), 404
                    
            except Exception as path_error:
                logger.error(f"Path validation error for {filename}: {str(path_error)}")
                return jsonify({'error': 'Invalid file path'}), 400
            
            # Check if file exists
            if not os.path.exists(file_path):
                logger.warning(f"File not found: {file_path}")
                return jsonify({'error': 'File not found'}), 404
            
            logger.info(f"Serving file: {filename} from {file_path}")
            
            # Serve the file using the directory and filename separately
            return send_from_directory(upload_folder, filename)
            
        except Exception as e:
            logger.error(f"Error serving file {filename}: {str(e)}")
            import traceback
            logger.error(f"Traceback: {traceback.format_exc()}")
            return jsonify({'error': 'Error serving file'}), 500
    
    @app.route('/file-stats/<user_id>', methods=['GET'])
    def get_file_stats(user_id):
        """Get file statistics for a user"""
        try:
            # Get file statistics
            pipeline = [
                {'$match': {'user_id': user_id, 'is_active': True}},
                {
                    '$group': {
                        '_id': None,
                        'total_files': {'$sum': 1},
                        'total_size': {'$sum': '$file_size'},
                        'file_types': {'$addToSet': '$file_extension'}
                    }
                }
            ]
            
            stats = list(files_collection.aggregate(pipeline))
            
            if stats:
                stat_data = stats[0]
                return jsonify({
                    'status': 'success',
                    'stats': {
                        'total_files': stat_data['total_files'],
                        'total_size': stat_data['total_size'],
                        'total_size_formatted': format_file_size(stat_data['total_size']),
                        'file_types': stat_data['file_types']
                    }
                })
            else:
                return jsonify({
                    'status': 'success',
                    'stats': {
                        'total_files': 0,
                        'total_size': 0,
                        'total_size_formatted': '0 B',
                        'file_types': []
                    }
                })
                
        except Exception as e:
            logger.error(f"Error retrieving file stats: {str(e)}")
            return jsonify({'error': 'Error retrieving file statistics'}), 500

    logger.info("File management routes registered successfully")
